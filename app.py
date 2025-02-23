import os
import uuid
import sqlite3
from pptx import Presentation
from datetime import datetime, timedelta
from pptx.util import Pt, Inches
import textwrap
from flask import Flask, render_template, request, jsonify, send_file, session
import json
from flask_cors import CORS
from openai import OpenAI
import uuid
from dotenv import load_dotenv

app = Flask(__name__)
app.secret_key = os.urandom(24)  # Required for using Flask sessions
CORS(app)  # Enables CORS for all routes by default

# Load environment variables from .env file
load_dotenv()
# Get the API key from environment variables
api_key = os.getenv("OPENAI_API_KEY")
print(api_key)

# Initialize OpenAI client
client = OpenAI(api_key=api_key)

# Function to connect to SQLite
def get_db_connection():
    conn = sqlite3.connect("chatbot_cache.db")
    conn.execute('''
        CREATE TABLE IF NOT EXISTS topic_explanations (
            course TEXT, week TEXT, day TEXT, topic TEXT, explanation TEXT,
            PRIMARY KEY (course, week, day, topic)
        )
    ''')
    conn.execute('''
        CREATE TABLE IF NOT EXISTS conversation_history (
            session_id TEXT, user_message TEXT, bot_response TEXT, timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.execute('''
        CREATE TABLE IF NOT EXISTS session_activity (
            session_id TEXT PRIMARY KEY, last_active DATETIME
        )
    ''')
    return conn

def reset_all_conversations():
    conn = get_db_connection()
    conn.execute("DELETE FROM conversation_history")
    conn.execute("DELETE FROM session_activity")
    conn.commit()
    conn.close()
    print("✅ All conversation buffers reset on server start.")

def update_session_activity(session_id):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        INSERT INTO session_activity (session_id, last_active) VALUES (?, ?)
        ON CONFLICT(session_id) DO UPDATE SET last_active=excluded.last_active
    ''', (session_id, datetime.utcnow()))
    conn.commit()
    conn.close()

def clear_inactive_sessions(timeout_minutes=30):
    cutoff_time = datetime.utcnow() - timedelta(minutes=timeout_minutes)
    conn = get_db_connection()
    cursor = conn.cursor()

    # Find inactive sessions
    cursor.execute('''
        SELECT session_id FROM session_activity WHERE last_active < ?
    ''', (cutoff_time,))
    inactive_sessions = cursor.fetchall()

    # Delete conversations and activity records
    for (inactive_session_id,) in inactive_sessions:
        cursor.execute('DELETE FROM conversation_history WHERE session_id = ?', (inactive_session_id,))
        cursor.execute('DELETE FROM session_activity WHERE session_id = ?', (inactive_session_id,))

    conn.commit()
    conn.close()

def save_conversation(session_id, user_message, bot_response):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        INSERT INTO conversation_history (session_id, user_message, bot_response) VALUES (?, ?, ?)
    ''', (session_id, user_message, bot_response))
    conn.commit()
    conn.close()

def get_conversation_history(session_id, limit=10):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        SELECT user_message, bot_response FROM conversation_history
        WHERE session_id = ? ORDER BY timestamp DESC LIMIT ?
    ''', (session_id, limit))
    history = cursor.fetchall()
    conn.close()
    return history[::-1]  # Return in chronological order


    
# Helper function to interact with GPT
def chat_with_gpt(prompt, type=None):
    try:
        if type:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                response_format={"type": "json_object"},
                messages=[{"role": "user", "content": prompt}],
            )
            return response.choices[0].message.content
        else:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
            )
            return response.choices[0].message.content
    except Exception as e:
        return f"Error: {e}"
    
def save_string_to_file(text, filename="output.txt"):
    """
    Save a string to a text file.
    
    Parameters:
    - text (str): The string to save.
    - filename (str): The name of the file (default: "output.txt").
    
    Returns:
    - str: Confirmation message with the file path.
    """
    try:
        with open(filename, "w", encoding="utf-8") as file:
            file.write(text)
        return f"String saved to {filename}"
    except Exception as e:
        return f"An error occurred: {e}"


# Load syllabus
def load_syllabus():
    with open("course/syllabus.json", "r") as file:
        return json.load(file)

def get_available_courses(syllabus):
    return list(syllabus.keys())

syllabus_json = load_syllabus()

@app.route('/')
def home():
    # Generate a new session ID if not already present
    if 'session_id' not in session:
        session['session_id'] = str(uuid.uuid4())
    return render_template("index.html", session_id=session['session_id'])

@app.route('/get_courses', methods=['GET'])
def get_courses():
    session_id = request.args.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    courses = get_available_courses(syllabus_json)
    return jsonify(courses=courses)

@app.route('/get_weeks', methods=['POST'])
def get_weeks():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    course = request.json.get("course")
    if course in syllabus_json:  # Check if the course exists in the syllabus JSON
        weeks = list(syllabus_json[course].keys())
        return jsonify(weeks=weeks)
    return jsonify(error="Invalid course selected."), 400

@app.route('/get_days', methods=['POST'])
def get_days():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    course = request.json.get("course")
    week = request.json.get("week")

    if not course or course not in syllabus_json:
        return jsonify(error="Invalid or missing course."), 400

    if not week or week not in syllabus_json[course]:
        return jsonify(error="Invalid or missing week."), 400

    days = list(syllabus_json[course][week].keys())
    return jsonify(days=days)

@app.route('/get_topics', methods=['POST'])
def get_topics():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    course = request.json.get("course")
    week = request.json.get("week")
    day = request.json.get("day")

    if not course or course not in syllabus_json:
        return jsonify(error="Invalid or missing course."), 400

    if not week or week not in syllabus_json[course]:
        return jsonify(error="Invalid or missing week."), 400

    if not day or day not in syllabus_json[course][week]:
        return jsonify(error="Invalid or missing day."), 400

    topics = syllabus_json[course][week][day].get("Topics", [])
    return jsonify(topics=topics)

@app.route('/teach_topic', methods=['POST'])
def teach_topic():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    course = request.json.get("course")
    week = request.json.get("week")
    day = request.json.get("day")
    topic = request.json.get("topic")

    conn = get_db_connection()
    cursor = conn.cursor()

    # Check if explanation exists in the database
    cursor.execute("SELECT explanation FROM topic_explanations WHERE course = ? AND week = ? AND day = ? AND topic = ?", 
                   (course, week, day, topic))
    row = cursor.fetchone()

    if row:  # If explanation is found, return cached response
        explanation = row[0]
        conn.close()
        return jsonify(explanation=explanation, cached=True)

    # Otherwise, generate explanation using GPT
    topic_file_path = f"course/lecture_notes/{week}/{day}.txt"
    
    if not os.path.exists(topic_file_path):
        return jsonify(error="Lecture notes file not found."), 404

    with open(topic_file_path, 'r', encoding='utf-8') as file:
        topic_content = file.read()

    prompt = f"By going through '''{topic_content}'''. Explain the '{topic}' in simple terms suitable for a beginner."
    explanation = chat_with_gpt(prompt)

    # Save the response to the database
    cursor.execute("INSERT INTO topic_explanations (course, week, day, topic, explanation) VALUES (?, ?, ?, ?, ?)", 
                   (course, week, day, topic, explanation))
    conn.commit()
    conn.close()

    return jsonify(explanation=explanation, cached=False)

@app.route('/chat', methods=['POST'])
def chat():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403

    message = request.json.get("message")
    if not message:
        return jsonify(error="Message is required."), 400

    # Clear inactive sessions and update current session
    clear_inactive_sessions(timeout_minutes=30)
    update_session_activity(session_id)

    history = get_conversation_history(session_id)
    history_prompt = "\n".join([f"User: {h[0]}\nBot: {h[1]}" for h in history])

    course = request.json.get("course")
    week = request.json.get("week")
    day = request.json.get("day")
    topic = request.json.get("topic")
    

    if not course or not week or not day or not topic:
        return jsonify(error="All selections (course, week, day, topic) are required to start the quiz."), 400

    topic_content = open(f"course/lecture_notes/{week}/{day}.txt", 'r', encoding='utf-8').read()
    
    if topic_content.strip() == "":
        return jsonify(error="No content found for the selected topic."), 404

    prompt = f"""
    Topic {topic} from the course {course} is about:
    {topic_content}
    Conversation History:
    {history_prompt}

    Current User Input: {message}

    Respond accordingly and don't go out of the scope of the conversation.
    """
    reply = chat_with_gpt(prompt)
    save_conversation(session_id, message, reply)

    return jsonify(reply=reply)


@app.route('/quiz')
def quiz_page():
    return render_template('quiz.html')  # Ensure 'quiz.html' is in the templates folder

@app.route('/start_quiz', methods=['POST'])
def start_quiz():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    course = request.json.get("course")
    week = request.json.get("week")
    day = request.json.get("day")
    topic = request.json.get("topic")

    if not course or not week or not day or not topic:
        return jsonify(error="All selections (course, week, day, topic) are required to start the quiz."), 400

    topic_content = open(f"course/lecture_notes/{week}/{day}.txt", 'r', encoding='utf-8').read()
    
    if topic_content.strip() == "":
        return jsonify(error="No content found for the selected topic."), 404

    prompt = f'''
                Create a short quiz with 3 multiple-choice questions and their answers
                on the topic "{topic_content}".
                Hint: the quiz should not be out of the topic.
                Present the quiz in JSON format as given here:
                {json.load(open("quiz_template.json"))}
                '''
    quiz_content = chat_with_gpt(prompt, type="json_object")

    try:
        quiz_data = json.loads(quiz_content)
    except json.JSONDecodeError:
        return jsonify(error="Failed to generate valid quiz data."), 500

    return jsonify(quiz=quiz_data)


@app.route('/submit_quiz', methods=['POST'])
def submit_quiz():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    responses = request.json.get("responses")
    quiz = request.json.get("quiz")

    if not responses or not quiz:
        return jsonify(error="Responses and quiz data are required."), 400

    score = 0
    total = len(quiz['questions'])
    feedback = []

    for i, response in enumerate(responses):
        correct_option_id = quiz['questions'][i]["correct_option_id"]
        if response == correct_option_id:
            score += 1
            feedback.append({"question": quiz['questions'][i]["question"], "correct": True})
        else:
            correct_answer_text = next(option["text"] for option in quiz['questions'][i]["options"] if option["id"] == correct_option_id)
            feedback.append({
                "question": quiz['questions'][i]["question"],
                "correct": False,
                "correct_answer": correct_answer_text
            })

    return jsonify(score=score, total=total, feedback=feedback)

# PowerPoint Generation Function
def create_visually_appealing_presentation(topic, slides_data, filename, max_chars_per_slide=1000):
    """
    Generate a visually appealing PowerPoint presentation with markdown content, 
    handling long text by splitting it across multiple slides and ensuring proper layout.
    """
    def add_content_slide(prs, title, content, layout):
        """
        Helper function to add a content slide to the presentation.
        Dynamically adjusts content to fit the slide.
        """
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        content_placeholder = slide.placeholders[1]

        # Dynamically adjust font size to fit the content
        font_size = 18  # Start with a standard font size
        max_lines = 12  # Max lines to fit in the placeholder
        content_lines = content.split("\n")

        # If content exceeds max lines, reduce font size iteratively
        while len(content_lines) > max_lines and font_size > 10:
            font_size -= 1
            max_lines += 2  # Slightly increase lines allowed as font size decreases

        # Add content to placeholder with adjusted font size
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear default content
        for line in content_lines:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.line_spacing = Pt(font_size + 6)  # Line spacing based on font size

    # Initialize the presentation
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Add title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = topic
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"An in-depth presentation on {topic}"

    # Customize title slide fonts
    title_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.italic = True

    # Add content slides
    for slide_data in slides_data:
        title = slide_data["title"]
        content = slide_data["content"]

        # Split content into chunks based on max_chars_per_slide
        content_chunks = textwrap.wrap(content, max_chars_per_slide, break_long_words=False, replace_whitespace=False)

        for i, chunk in enumerate(content_chunks):
            slide_title = title if i == 0 else f"{title} (Cont.)"
            add_content_slide(prs, slide_title, chunk, content_slide_layout)

    # Save the presentation
    try:
        prs.save(filename)
    except Exception as e:
        raise IOError(f"Failed to save the presentation: {e}")


@app.route('/generate_presentation', methods=['POST'])
def generate_presentation():
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    data = request.json
    informational_text = data.get("informational_text")
    topic = data.get("topic", "Presentation")

    if not informational_text or len(informational_text.strip()) < 10:
        return jsonify(error="Informational text is required and must be meaningful."), 400

    try:
        # Generate slide titles
        titles_prompt = f"Don't add preamble and Based on the following informational text, generate list of slide titles for a presentation separated by '\\n':\n\n{informational_text}"
        titles_response = client.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": titles_prompt}]
        )
        titles_content = titles_response.choices[0].message.content if titles_response.choices and titles_response.choices[0].message.content else None

        if not titles_content:
            return jsonify(error="Failed to generate slide titles. OpenAI response was empty or invalid."), 500

        slide_titles = [line.strip() for line in titles_content.split("\n") if line.strip()]

        if not slide_titles:
            return jsonify(error="No slide titles generated. Please review the informational text."), 400

        # Generate slide content
        slides_data = []
        for title in slide_titles:
            content_prompt = f"Don't add preamble. Generate short content for the presentation slide titled '{title}' based on the following text:\n\n{informational_text}"
            content_response = client.chat.completions.create(
                model="gpt-4",
                messages=[{"role": "user", "content": content_prompt}]
            )
            content = content_response.choices[0].message.content if content_response.choices and content_response.choices[0].message.content else None

            if not content:
                return jsonify(error=f"Content for slide '{title}' is empty or invalid."), 500

            slides_data.append({
                "title": title,
                "content": content.strip()
            })

        # Generate the presentation
        output_dir = "generated_ppt"
        os.makedirs(output_dir, exist_ok=True)
        unique_id = uuid.uuid4().hex
        ppt_filename = os.path.join(output_dir, f"presentation_{unique_id}.pptx")
        create_visually_appealing_presentation(topic, slides_data, ppt_filename)

        return jsonify(message="Presentation generated successfully!", download_url=f"/download_ppt?filename=presentation_{unique_id}.pptx")
    except Exception as e:
        return jsonify(error=f"Failed to generate presentation: {str(e)}"), 500


@app.route('/download_ppt', methods=['GET'])
def download_ppt():
    """
    Download the generated PowerPoint presentation.
    """
    session_id = request.json.get("session_id")
    if session_id != session.get('session_id'):
        return jsonify(error="Invalid session."), 403
    
    filename = request.args.get("filename")
    if not filename:
        return jsonify(error="Filename is required."), 400

    filepath = os.path.join("generated_ppt", filename)
    if not os.path.exists(filepath):
        return jsonify(error="File not found."), 404

    return send_file(filepath, as_attachment=True)

if __name__ == "__main__":
    reset_all_conversations()
    app.run(debug=True)